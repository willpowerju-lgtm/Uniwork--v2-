#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
registry_check.py — 提交前的「registry 一致性闸」。

用途：focus 标的（Vault/Company Coverage/<T>/ 下有 _data_registry.json）下游派生件
（研报 docx / wiki md / 导出 xlsx·pptx）被改动并 commit 时，核对里面的"数字"是否还
和 _data_registry.json（SSOT，由 model 派生）对得上。

核心思想（diff-from-HEAD，精准、低误报）：
  只盯"用户这次动过的数字"。对每条 distinctive 的 registry 数值 V：
    - 它原来出现在文件的 HEAD 版本里(in_old) 但改完后不见了(not in_new) → 判定为
      "下游被手改、和 model 对不上" 的 divergence。
  纯小数位/取整变化（1761→1761.0、36.46→36.5、9→9.0）经 rounding-equiv 吸收，不算 divergence。

只核对 source=="Excel model" 的 model 派生条目（assumptions/forecast），且只取
distinctive 值（|v|>=50 或带小数，且非 1990-2099 年份整数），避免小整数/日期误报。

退出码：有 divergence → 1（当 gate）；一致或无基线 → 0。stdout 永远是一段 JSON。
"""
import sys, io, os, re, json, argparse
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')

# 数字 token：可选 $、千分位逗号、小数、可选 %；括号记账负数另算
NUM_RE = re.compile(r'\(?-?\$?\s?\d[\d,]*(?:\.\d+)?\s?%?\)?')

def _norm_token(s):
    """'(1,761)'→(-1761.0,0) ; '36.5%'→(36.5,1) ; '$1,761.0'→(1761.0,1)。返回 (value, decimals) 或 None"""
    raw = s.strip()
    neg = raw.startswith('(') and raw.endswith(')')
    t = raw.strip('()').replace(',', '').replace('$', '').replace('%', '').replace(' ', '')
    if t in ('', '-', '.', '-.'):
        return None
    try:
        v = float(t)
    except ValueError:
        return None
    if neg:
        v = -v
    dec = len(t.split('.')[1]) if '.' in t else 0
    return (v, dec)

def extract_tokens_from_text(text):
    out = []
    for m in NUM_RE.findall(text or ''):
        nt = _norm_token(m)
        if nt is not None:
            out.append(nt)
    return out

def extract_tokens(path):
    """按文件类型抽出所有数字 token。返回 [(value, decimals), ...]。读不了就返回 None。"""
    if not path or not os.path.exists(path):
        return None
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext == '.docx':
            import docx
            d = docx.Document(path)
            parts = [p.text for p in d.paragraphs]
            for tb in d.tables:
                for row in tb.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            return extract_tokens_from_text('\n'.join(parts))
        if ext == '.xlsx':
            import openpyxl
            wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
            toks = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    for v in row:
                        if isinstance(v, bool):
                            continue
                        if isinstance(v, (int, float)):
                            dec = len(str(v).split('.')[1]) if ('.' in str(v)) else 0
                            toks.append((float(v), dec))
                        elif isinstance(v, str):
                            toks.extend(extract_tokens_from_text(v))
            wb.close()
            return toks
        if ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(path)
            parts = []
            for sl in prs.slides:
                for sh in sl.shapes:
                    if sh.has_text_frame:
                        parts.append(sh.text_frame.text)
                    if sh.has_table:
                        for r in sh.table.rows:
                            for c in r.cells:
                                parts.append(c.text)
            return extract_tokens_from_text('\n'.join(parts))
        # .md / .txt / .json / .csv / 其它文本
        with open(path, 'r', encoding='utf-8', errors='replace') as f:
            return extract_tokens_from_text(f.read())
    except Exception as e:
        print(json.dumps({'ok': True, 'skipped': True, 'reason': f'读取失败:{type(e).__name__}:{e}'}, ensure_ascii=False))
        sys.exit(0)

def rounding_equiv(V, T, dT):
    """报告里显示的 T（dT 位小数）是不是 registry 真值 V 的取整/小数位版本。"""
    if round(V, dT) == round(T, dT):
        return True
    return abs(T - V) / max(abs(V), 1e-9) <= 0.005   # 0.5% 兜底，吸收显示舍入

def present(V, tokens):
    return any(rounding_equiv(V, T, dT) for (T, dT) in tokens)

def distinctive(v):
    """够"特异"才核对：|v|>=50 或带小数；排除 1990-2099 年份整数（满天飞的日期）。"""
    if float(v).is_integer() and 1990 <= v <= 2099:
        return False
    return abs(v) >= 50 or (v != round(v))

def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('--registry', required=True)
    ap.add_argument('--new', required=True, help='改动后的文件（盘上当前内容）')
    ap.add_argument('--old', default=None, help='HEAD 版本（server 用 git show 落到临时文件）')
    ap.add_argument('--label', default=None, help='展示用文件名/相对路径')
    a = ap.parse_args()

    label = a.label or os.path.basename(a.new)
    try:
        reg = json.load(open(a.registry, encoding='utf-8'))
    except Exception as e:
        print(json.dumps({'ok': True, 'skipped': True, 'reason': f'registry 读不了:{e}'}, ensure_ascii=False)); return 0

    entries = reg.get('entries', {})
    # 只核对 model 派生、distinctive 的数值条目
    checks = []
    for k, ent in entries.items():
        if not isinstance(ent, dict):
            continue
        v = ent.get('value')
        if not isinstance(v, (int, float)) or isinstance(v, bool):
            continue
        if (ent.get('source') or '') != 'Excel model':
            continue
        if not distinctive(float(v)):
            continue
        checks.append((k, float(v), ent.get('unit'), ent.get('period')))

    new_tokens = extract_tokens(a.new)
    if new_tokens is None:
        print(json.dumps({'ok': True, 'skipped': True, 'reason': 'new 文件读不了'}, ensure_ascii=False)); return 0

    old_tokens = extract_tokens(a.old) if a.old else None
    if old_tokens is None:
        # 无基线（首次提交 / 二进制 HEAD 取不到）→ 无法判定"改了什么"，不拦
        print(json.dumps({'ok': True, 'mode': 'no-baseline', 'file': label,
                          'checked_entries': len(checks),
                          'note': '无 HEAD 基线，跳过 diff（首次提交或取不到旧版本）'}, ensure_ascii=False))
        return 0

    divergences = []
    for k, V, unit, period in checks:
        in_old = present(V, old_tokens)
        in_new = present(V, new_tokens)
        if in_old and not in_new:
            # 这条 model 数字原来在报告里、改完不见了 → 找最接近的新数字当"疑似被替换成"
            nearest = min(new_tokens, key=lambda t: abs(t[0] - V), default=None)
            divergences.append({
                'key': k, 'period': period, 'unit': unit,
                'registry_value': V,
                'gone_replaced_by': (nearest[0] if nearest else None),
            })

    ok = not divergences
    print(json.dumps({
        'ok': ok,
        'file': label,
        'checked_entries': len(checks),
        'divergences': divergences,
        'note': ('一致：本次改动未动到 model 派生数字' if ok else
                 f'⚠️ 检测到 {len(divergences)} 个 model 派生数字被改/移除，下游与 model 对不上'),
    }, ensure_ascii=False))
    return 0 if ok else 1

if __name__ == '__main__':
    sys.exit(main())
